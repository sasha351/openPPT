"""
title: Export to PowerPoint
author: openPPT
version: 0.7.3
requirements: python-pptx
description: Adds an "Export to PowerPoint" button that converts an HTML outline in the assistant message into a downloadable .pptx, with tables, code blocks, speaker notes, images, and custom templates. Attach a .pptx/.potx to the chat and the deck inherits its theme, fonts, and layouts — the model just dumps content as an outline and picks layouts by name. Works with any model (no tool calling needed).
"""

import html
import inspect
import io
import re
import urllib.request
import uuid
from html.parser import HTMLParser

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Precedes anything openPPT appends to a chat message — the download link, the
# no-outline diagnostic. An HTML comment renders invisibly in chat but survives
# a round trip back in, and the parser stops dead at this comment (see
# _OutlineParser.handle_comment), so re-clicking Export can't parse openPPT's
# own output back into slides (v0.3.5 exported its own error report as a deck).
APPENDIX_MARKER = "<!-- openppt -->"

# Kept equal to the docstring's 'version:' line — a stale paste in Open WebUI's
# Functions list is otherwise invisible, and the diagnostic is where it shows.
VERSION = "0.7.3"


def _new_slide(title: str, layout: str = "") -> dict:
    return {
        "title": title,
        "subtitle": "",
        "layout": layout,
        "notes": "",
        "bullets": [],
        "table": [],
        "code": "",
    }


def _attr(attrs: dict, name: str) -> str:
    return (attrs.get(name) or "").strip()


class _OutlineParser(HTMLParser):
    """Turns a stream of tag events into the same slide-dict shape the old
    Markdown parser produced, so build_pptx and everything downstream needs
    no changes.

    Deliberately tolerant rather than a strict-HTML validator, because local
    models don't emit strict HTML any more reliably than they emitted strict
    Markdown: unclosed '<slide>'/'<li>'/'<p>' tags are auto-closed, stray text
    directly inside '<slide>' (not wrapped in '<p>'/'<li>') still becomes a
    bullet instead of being dropped, and '<blockquote>' is accepted as a
    synonym for '<notes>'. HTMLParser itself is already lenient about
    malformed markup (mismatched tags, missing quotes, bad nesting) and never
    raises on it, which is most of why HTML beats hand-rolled regex here.
    """

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.slides = []
        self.current = None
        self.in_slide = False
        self.stopped = False
        self._reset_slide_state()

    def _reset_slide_state(self):
        self.list_depth = 0
        self.li_stack = []  # [{"level", "buf": [str], "image": dict|None, "emitted": bool}]
        self.pre_depth = 0
        self.code_buf = []
        self.notes_depth = 0
        self.notes_buf = []
        self.p_depth = 0
        self.p_buf = []
        self.table = None
        self.row = None
        self.cell_depth = 0
        self.cell_buf = []

    def _close_slide(self):
        """Flush whatever was still open (unclosed li/p/notes/table/pre) and
        end the current slide — called on '</slide>', on a new '<slide>'
        opening without a previous close, and at end-of-document."""
        while self.li_stack:
            frame = self.li_stack.pop()
            if not frame["emitted"] and frame["image"] is not None:
                self.current["bullets"].append((frame["level"], frame["image"]))
            text = "".join(frame["buf"]).strip()
            if text:
                self.current["bullets"].append((frame["level"], text))
        if self.p_depth > 0:
            text = "".join(self.p_buf).strip()
            if text:
                self.current["bullets"].append((0, text))
        if self.notes_depth > 0:
            text = "".join(self.notes_buf).strip()
            if text:
                self.current["notes"] = (self.current["notes"] + "\n" + text).strip()
        if self.table is not None:
            if self.row is not None:
                self.table.append(self.row)
            self.current["table"].extend(self.table)
        if self.pre_depth > 0:
            self.current["code"] += "".join(self.code_buf)
        self._reset_slide_state()
        self.in_slide = False
        self.current = None

    def handle_comment(self, data):
        if data.strip() == "openppt":
            self.stopped = True  # hard stop: never parse openPPT's own appendix

    def handle_starttag(self, tag, attrs):
        if self.stopped:
            return
        a = dict(attrs)
        if self.pre_depth > 0 and tag != "pre":
            if tag != "code":  # transparent — doesn't count as nesting depth
                self.code_buf.append(self.get_starttag_text())
            return
        if tag == "pre":
            self.pre_depth += 1
            return
        if tag == "slide":
            if self.in_slide:
                self._close_slide()  # tolerate a missing '</slide>'
            self.current = _new_slide(_attr(a, "title"), _attr(a, "layout"))
            self.current["subtitle"] = _attr(a, "subtitle")
            self.slides.append(self.current)
            self.in_slide = True
            return
        if not self.in_slide:
            return  # stray tag before the first '<slide>' — ignore
        if tag in ("ul", "ol"):
            if self.li_stack:
                # A nested list inside a still-open '<li>' — flush that li's
                # own text now, in document order, rather than waiting for its
                # '</li>' (which comes after the nested list closes and would
                # put the parent's text after its children's).
                top = self.li_stack[-1]
                if not top["emitted"]:
                    if top["image"] is not None:
                        self.current["bullets"].append((top["level"], top["image"]))
                    else:
                        text = "".join(top["buf"]).strip()
                        if text:
                            self.current["bullets"].append((top["level"], text))
                    top["buf"] = []
                    top["emitted"] = True
            self.list_depth += 1
        elif tag == "li":
            level = min(self.list_depth - 1, 4) if self.list_depth > 0 else 0
            self.li_stack.append(
                {"level": max(level, 0), "buf": [], "image": None, "emitted": False}
            )
        elif tag == "img":
            image = {"alt": _attr(a, "alt"), "image": _attr(a, "src")}
            if self.li_stack:
                self.li_stack[-1]["image"] = image
            else:
                self.current["bullets"].append((0, image))
        elif tag in ("notes", "blockquote"):
            self.notes_depth += 1
        elif tag == "p":
            self.p_depth += 1
        elif tag == "table":
            self.table = []
        elif tag == "tr":
            self.row = []
        elif tag in ("td", "th"):
            self.cell_depth += 1

    def handle_endtag(self, tag):
        if self.stopped:
            return
        if self.pre_depth > 0 and tag != "pre":
            if tag != "code":
                self.code_buf.append(f"</{tag}>")
            return
        if tag == "pre":
            self.pre_depth = max(0, self.pre_depth - 1)
            if self.pre_depth == 0 and self.in_slide:
                self.current["code"] += "".join(self.code_buf)
                self.code_buf = []
            return
        if tag == "slide":
            if self.in_slide:
                self._close_slide()
            return
        if not self.in_slide:
            return
        if tag in ("ul", "ol"):
            self.list_depth = max(0, self.list_depth - 1)
        elif tag == "li":
            if self.li_stack:
                frame = self.li_stack.pop()
                if not frame["emitted"] and frame["image"] is not None:
                    self.current["bullets"].append((frame["level"], frame["image"]))
                text = "".join(frame["buf"]).strip()
                if text:
                    self.current["bullets"].append((frame["level"], text))
        elif tag in ("notes", "blockquote"):
            self.notes_depth = max(0, self.notes_depth - 1)
            if self.notes_depth == 0:
                text = "".join(self.notes_buf).strip()
                self.notes_buf = []
                if text:
                    self.current["notes"] = (self.current["notes"] + "\n" + text).strip()
        elif tag == "p":
            self.p_depth = max(0, self.p_depth - 1)
            if self.p_depth == 0:
                text = "".join(self.p_buf).strip()
                self.p_buf = []
                if text:
                    self.current["bullets"].append((0, text))
        elif tag == "table":
            if self.table is not None:
                self.current["table"].extend(self.table)
                self.table = None
        elif tag == "tr":
            if self.row is not None and self.table is not None:
                self.table.append(self.row)
            self.row = None
        elif tag in ("td", "th"):
            if self.cell_depth > 0:
                self.cell_depth -= 1
                text = "".join(self.cell_buf).strip()
                self.cell_buf = []
                if self.row is not None:
                    self.row.append(text)

    def handle_data(self, data):
        if self.stopped or not self.in_slide:
            return
        if self.pre_depth > 0:
            self.code_buf.append(data)
        elif self.cell_depth > 0:
            self.cell_buf.append(data)
        elif self.li_stack:
            self.li_stack[-1]["buf"].append(data)
        elif self.p_depth > 0:
            self.p_buf.append(data)
        elif self.notes_depth > 0:
            self.notes_buf.append(data)
        else:
            # Bare text directly inside '<slide>', not wrapped in '<p>'/'<li>'
            # — models forget the wrapper constantly, so treat it as a bullet
            # rather than silently dropping the content.
            text = data.strip()
            if text:
                self.current["bullets"].append((0, text))


def _extract_fences(text: str) -> str:
    """Handle ``` fences the model adds despite being asked for plain HTML.

    A fence opened before the first '<slide' is a wrapper around the whole
    answer (models do this constantly) — unwrapped and parsed through. A
    fence opened after is a real code sample; its content is HTML-escaped and
    spliced in as a '<pre>' block so the outline parser needs only one code
    path ('<pre>', wherever the text came from), and code containing '<'/'>'
    can't be mistaken for markup.
    """
    seen_slide = False
    in_fence = False
    fence_is_code = False
    code_lines = []
    out = []
    for raw in text.splitlines():
        if raw.strip().startswith("```"):
            if not in_fence:
                in_fence, fence_is_code, code_lines = True, seen_slide, []
            else:
                in_fence = False
                if fence_is_code:
                    out.append(f"<pre>{html.escape(chr(10).join(code_lines))}</pre>")
                else:
                    out.extend(code_lines)
            continue
        if in_fence:
            code_lines.append(raw)
            continue
        if "<slide" in raw.lower():
            seen_slide = True
        out.append(raw)
    if in_fence:  # unclosed fence at EOF — flush what we have instead of losing it
        if fence_is_code:
            out.append(f"<pre>{html.escape(chr(10).join(code_lines))}</pre>")
        else:
            out.extend(code_lines)
    return "\n".join(out)

# Placeholder type groups (see PP_PLACEHOLDER).
TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
SUBTITLE_TYPES = (PP_PLACEHOLDER.SUBTITLE,)
BODY_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
# Placeholders that are chrome, never content we want to fill.
CHROME_TYPES = (
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.HEADER,
)


def parse_outline(markup: str) -> list:
    """Parse an HTML outline into slides.

    Each slide is {"title", "subtitle", "layout", "notes", "table", "code",
    "bullets": [(level, text) | (level, {"image": url, "alt": alt})]}.

    Each '<slide title="..." subtitle="..." layout="...">' element is one
    slide; 'subtitle' and 'layout' are optional, and 'layout' names a
    template layout matched against the template's layouts in build_pptx.
    Inside a slide: '<ul>'/'<ol>' of '<li>' are bullets, with a nested
    '<ul>'/'<ol>' inside an '<li>' nesting the bullet under it; '<p>' is a
    single bullet; '<notes>' (or '<blockquote>') is appended to the slide's
    speaker notes; '<img src="..." alt="...">' — bare or inside an '<li>' —
    embeds an image instead of text; '<table>' of '<tr>' of '<td>'/'<th>'
    becomes a real table; and '<pre>' (with or without a nested '<code>')
    becomes a code block rendered in a monospace box. A ``` fence is also
    accepted as a code block for models that fall back to it out of habit —
    escaped and folded into an equivalent '<pre>' before parsing.

    The parser (html.parser.HTMLParser, stdlib) is deliberately tolerant
    rather than a strict validator: an unclosed '<slide>'/'<li>'/'<p>' is
    auto-closed by whatever ends it (the next '<slide>', or end of input),
    and bare text sitting directly inside '<slide>' — not wrapped in '<p>' or
    '<li>' — still becomes a bullet instead of being silently dropped. Text
    outside any '<slide>' (chat preamble, unknown tags) is ignored. A ```
    fence that opens before the first '<slide>' is a wrapper around the whole
    answer, not a code sample, and is unwrapped rather than swallowing the
    deck — models wrap their answer in one whether or not you ask them to.

    Returns [] if no slides found. Parsing stops dead at the first
    `APPENDIX_MARKER` HTML comment, so re-exporting a message openPPT already
    appended to (a download link, a diagnostic) never parses that text back
    into slides — even if the model's own last '<slide>' was left unclosed.
    """
    parser = _OutlineParser()
    try:
        parser.feed(_extract_fences(markup))
        parser.close()
    except Exception:
        pass  # never raise — whatever slides were parsed before the failure still count
    if parser.in_slide:
        parser._close_slide()
    return parser.slides


async def _maybe_await(value):
    """Open WebUI's Files API turned async; a pasted Function must fit both.

    Calling an async Files method without awaiting returns a coroutine that
    never runs — the insert silently doesn't happen and the download 404s.
    """
    return await value if inspect.isawaitable(value) else value


def _fetch_image(url: str):
    """Return image bytes for a local path or http(s) URL, or None on failure."""
    try:
        if url.startswith("http://") or url.startswith("https://"):
            with urllib.request.urlopen(url, timeout=10) as resp:
                return io.BytesIO(resp.read())
        with open(url, "rb") as f:
            return io.BytesIO(f.read())
    except Exception:
        return None


def _file_entries(files):
    """(id, filename) for each attachment entry. Open WebUI has shipped both
    a flat {'id','name'} and a nested {'file': {'id','filename'}} shape."""
    for f in files or []:
        if not isinstance(f, dict):
            continue
        inner = f.get("file") if isinstance(f.get("file"), dict) else {}
        fid = inner.get("id") or f.get("id")
        if fid:
            yield fid, (f.get("name") or inner.get("filename") or "")


def _ph_type(placeholder):
    """Placeholder type, tolerant of odd templates that leave it unset."""
    try:
        return placeholder.placeholder_format.type
    except Exception:
        return None


def _clear_slides(prs) -> None:
    """Remove any example slides from a template so the deck starts clean.

    A .potx has none, but a .pptx used as a template often carries sample
    slides; the template's masters, layouts, and theme are untouched.
    """
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rid:
            prs.part.drop_rel(rid)
        sld_id_lst.remove(sld_id)


def _layout_stats(layout):
    """Summarise a layout by placeholder role for role-based selection."""
    titles = subtitles = bodies = 0
    for ph in layout.placeholders:
        t = _ph_type(ph)
        if t in TITLE_TYPES:
            titles += 1
        elif t in SUBTITLE_TYPES:
            subtitles += 1
        elif t in BODY_TYPES:
            bodies += 1
    return {"title": titles, "subtitle": subtitles, "body": bodies}


def list_layouts(prs_or_bytes) -> list:
    """Return [(name, {'title','subtitle','body' counts})] for a template.

    Accepts a Presentation, template bytes, or a file-like/stream. Handy for
    telling the model which layout names it can pick from.
    """
    prs = _as_presentation(prs_or_bytes)
    return [(lay.name, _layout_stats(lay)) for lay in prs.slide_layouts]


def _as_presentation(template):
    """Coerce None / bytes / stream / Presentation into a Presentation."""
    if template is None:
        return Presentation()
    if hasattr(template, "slide_layouts"):
        return template
    if isinstance(template, (bytes, bytearray)):
        return Presentation(io.BytesIO(template))
    return Presentation(template)  # file path or file-like


def _find_layout_by_name(prs, name: str):
    """Case-insensitive match of a layout by name (exact, then substring)."""
    if not name:
        return None
    target = name.strip().lower()
    layouts = list(prs.slide_layouts)
    for lay in layouts:
        if lay.name.lower() == target:
            return lay
    for lay in layouts:
        if target in lay.name.lower() or lay.name.lower() in target:
            return lay
    return None


def _pick_layout(prs, role: str):
    """Choose a layout for a role ('title', 'content', 'title_only') by
    inspecting placeholders, so it works on any template regardless of the
    order or names its layouts happen to use."""
    layouts = list(prs.slide_layouts)
    scored = [(lay, _layout_stats(lay)) for lay in layouts]

    if role == "title":
        # Prefer title + subtitle, no bullet body (a real title slide).
        for lay, s in scored:
            if s["title"] and s["subtitle"] and not s["body"]:
                return lay
        for lay, s in scored:
            if s["title"] and s["subtitle"]:
                return lay
    if role == "title_only":
        # A title, nothing else to fill.
        for lay, s in scored:
            if s["title"] and not s["subtitle"] and not s["body"]:
                return lay
    # role == "content" (and fallbacks for the above): title + exactly one body.
    for lay, s in scored:
        if s["title"] and s["body"] == 1 and not s["subtitle"]:
            return lay
    for lay, s in scored:
        if s["title"] and s["body"]:
            return lay
    for lay, s in scored:
        if s["title"]:
            return lay
    return layouts[0]


def _title_placeholder(slide):
    if slide.shapes.title is not None:
        return slide.shapes.title
    for ph in slide.placeholders:
        if _ph_type(ph) in TITLE_TYPES:
            return ph
    return None


def _subtitle_placeholder(slide):
    for ph in slide.placeholders:
        if _ph_type(ph) in SUBTITLE_TYPES:
            return ph
    return None


def _body_placeholder(slide):
    """First fillable body placeholder that isn't the title/subtitle/chrome."""
    title = _title_placeholder(slide)
    title_id = id(title) if title is not None else None
    for ph in slide.placeholders:
        if id(ph) == title_id:
            continue
        t = _ph_type(ph)
        if t in BODY_TYPES:
            return ph
    # Fall back to any non-title, non-chrome, non-subtitle placeholder.
    for ph in slide.placeholders:
        if id(ph) == title_id:
            continue
        t = _ph_type(ph)
        if t in CHROME_TYPES or t in SUBTITLE_TYPES:
            continue
        return ph
    return None


def _title_size(text: str):
    """Point size for a title long enough to overflow, or None to keep the
    template's. python-pptx never recomputes PowerPoint's autofit, so without
    this a 200-character title renders at full size across the whole slide."""
    for limit, size in ((50, None), (80, 32), (120, 26), (200, 20)):
        if len(text) <= limit:
            return size
    return 16


def _body_size(bullets) -> int:
    """Point size for a bullet list, shrinking as it gets denser. Long bullets
    wrap, so count the lines they'll take rather than the bullets."""
    lines = sum(1 + len(text) // 60 for _, text in bullets)
    for limit, size in ((5, 20), (8, 16), (12, 14)):
        if lines <= limit:
            return size
    return 12


def _fit(text_frame) -> None:
    """Wrap, and set autofit so tools that recompute it (LibreOffice, Google
    Slides) can still shrink text further. PowerPoint itself only recomputes
    autofit on edit, so on first open our explicit Pt() sizes are what fit
    the text — this just avoids overflow markers on the placeholder."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_title(slide, text: str) -> None:
    ph = _title_placeholder(slide)
    if ph is None or not ph.has_text_frame:
        return
    ph.text_frame.text = text
    _fit(ph.text_frame)
    size = _title_size(text)
    if size:
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(size)


def _fill_bullets(text_frame, bullets) -> None:
    _fit(text_frame)
    size = Pt(_body_size(bullets))
    for j, (level, text) in enumerate(bullets):
        p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = size


def _content_rect(prs, slide):
    """Box for extra content (table, code): the body placeholder's, or a
    margin-inset box under the title when the layout has none."""
    body = _body_placeholder(slide)
    if body is not None:
        return body, (body.left, body.top, body.width, body.height)
    left = Emu(int(prs.slide_width * 0.08))
    top = Emu(int(prs.slide_height * 0.28))
    return None, (left, top, prs.slide_width - 2 * left, prs.slide_height - top - left)


def _drop_shape(shape) -> None:
    """Remove a shape (an unfilled placeholder renders as an empty prompt box)."""
    shape._element.getparent().remove(shape._element)


def _add_table(slide, rows, rect) -> None:
    left, top, width, height = rect
    cols = max(len(r) for r in rows)
    table = slide.shapes.add_table(len(rows), cols, left, top, width, height).table
    size = Pt(14 if len(rows) <= 6 else 11)
    for r, row in enumerate(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = size


def _add_code(slide, code, rect) -> None:
    lines = code.rstrip().splitlines()
    box = slide.shapes.add_textbox(*rect)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4)
    tf = box.text_frame
    tf.word_wrap = True
    # Monospace and small: code lines don't wrap gracefully, so err narrow.
    size = Pt(14 if len(lines) <= 10 else (11 if len(lines) <= 18 else 9))
    for j, text in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Consolas"
        p.font.size = size


def _split_rect(rect, fraction):
    """Top `fraction` of a box, and the remainder below it, with a gutter."""
    left, top, width, height = rect
    gutter = Pt(8)
    upper = int(height * fraction)
    return (left, top, width, Emu(upper)), (
        left,
        Emu(top + upper + gutter),
        width,
        Emu(height - upper - gutter),
    )


def build_pptx(slides: list, template=None) -> bytes:
    """Build a .pptx from parsed slides.

    template: optional .pptx/.potx as bytes, a stream, a path, or a
    Presentation. When given, the deck inherits that template's theme,
    fonts, and layouts, and each slide's 'layout' name is matched against the
    template's layouts. When omitted, the python-pptx default template is used.
    """
    prs = _as_presentation(template)
    _clear_slides(prs)

    for i, s in enumerate(slides):
        text_bullets = [(lvl, t) for lvl, t in s["bullets"] if isinstance(t, str)]
        image_bullets = [(lvl, t) for lvl, t in s["bullets"] if isinstance(t, dict)]
        table, code = s.get("table") or [], (s.get("code") or "").strip()
        image_only = not text_bullets and not table and not code and image_bullets
        is_title = i == 0 and not s["bullets"] and not table and not code

        layout = _find_layout_by_name(prs, s.get("layout", ""))
        if layout is None:
            role = "title" if is_title else ("title_only" if image_only else "content")
            layout = _pick_layout(prs, role)

        ps = prs.slides.add_slide(layout)
        _set_title(ps, s["title"])

        if s["subtitle"]:
            sub = _subtitle_placeholder(ps) or _body_placeholder(ps)
            if sub is not None:
                sub.text = s["subtitle"]

        body, rect = _content_rect(prs, ps)
        if text_bullets and body is not None and body.has_text_frame:
            _fill_bullets(body.text_frame, text_bullets)
            if table or code:
                # Bullets keep the top of the box, the table/code takes the rest.
                top_rect, rect = _split_rect(rect, 0.45)
                body.height = top_rect[3]
        elif (table or code) and body is not None:
            _drop_shape(body)  # nothing to put in it; don't render an empty box
        if table and code:
            table_rect, rect = _split_rect(rect, 0.5)
        else:
            table_rect = rect
        if table:
            _add_table(ps, table, table_rect)
        if code:
            _add_code(ps, code, rect)

        for _, img in image_bullets:
            data = _fetch_image(img["image"])
            if data is not None:
                ps.shapes.add_picture(data, Pt(60), Pt(120), height=Pt(340))

        if s["notes"]:
            ps.notes_slide.notes_text_frame.text = s["notes"]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _chat_messages(chat) -> list:
    """Messages out of a saved chat dict, oldest-first.

    Open WebUI keeps the same conversation twice — a flat 'messages' list and
    a 'history.messages' id->message map — and which one is populated varies
    by version, so read whichever holds more. Not "flat if it's non-empty":
    the flat list can lag a turn behind the history map (it's rebuilt from the
    current branch on save), and a flat list holding just the user's message
    would then shadow a history that has the assistant's reply in it.
    """
    if not isinstance(chat, dict):
        return []
    flat = chat.get("messages")
    flat = flat if isinstance(flat, list) else []
    history = chat.get("history")
    hist = history.get("messages") if isinstance(history, dict) else None
    if isinstance(hist, dict):
        hist = list(hist.values())
    elif not isinstance(hist, list):
        hist = []
    return hist if len(hist) > len(flat) else flat


def _shape(messages: list) -> str:
    """'user:str(21), assistant:list(2)' — role and content shape per message.

    Read only when no outline was found, and the four ways that happens need
    different fixes: content arriving as a list (multimodal), as an empty
    string, under no assistant role at all, or already consumed by the
    appendix strip. The character count alone can't tell them apart.
    """
    parts = []
    for msg in (messages or [])[-4:]:
        if not isinstance(msg, dict):
            parts.append(f"{type(msg).__name__}!")
            continue
        body = msg.get("content")
        size = len(body) if hasattr(body, "__len__") else "?"
        parts.append(f"{msg.get('role', '?')}:{type(body).__name__}({size})")
    return ", ".join(parts) or "(none)"


def _saved_shape(chat) -> str:
    """'flat=1 hist=3' — how many messages each of the saved chat's two message
    stores holds, or why there's no saved copy at all.

    Read only when the request carried no outline and the saved-chat fallback
    didn't rescue it, where three very different failures look identical from
    the outside: the chat lookup failed (Open WebUI API drift, temporary chat),
    the chat is saved but the assistant turn isn't in it yet, or both stores
    disagree and _chat_messages read the emptier one.
    """
    if chat is None:
        return "none"
    if not isinstance(chat, dict):
        return f"{type(chat).__name__}!"

    def size(v):
        return len(v) if hasattr(v, "__len__") else "?"

    history = chat.get("history")
    hist = history.get("messages") if isinstance(history, dict) else None
    return f"flat={size(chat.get('messages'))} hist={size(hist)}"


def _pick_outline(messages: list, target=None) -> str:
    """Content to export: the clicked assistant message, else the newest one
    that actually holds an outline.

    Open WebUI sometimes hands the action a message whose content is empty, and
    a strict clicked-message-only read then reports "no outline" at a chat that
    plainly has one. Returns the clicked/newest text when nothing parses, so
    the caller can echo what it actually read.
    """
    candidates = []
    for msg in reversed(messages or []):
        if msg.get("role") != "assistant":
            continue
        # Drop our own appended output (download link, diagnostic) so a
        # re-click echoes back what the model wrote, not what we last posted.
        # The blank line right before the marker is ours too (see the emits
        # below), so trim it along with everything after — but only when the
        # marker is actually present, or an ordinary message's own trailing
        # whitespace would get eaten.
        text = msg.get("content") or ""
        if not isinstance(text, str):
            # Open WebUI can hand back list-shaped (multimodal) content; treat
            # it as no outline rather than raising out of .partition() below.
            text = ""
        before, marker, _ = text.partition(APPENDIX_MARKER)
        if marker:
            text = before.rstrip()
        if target is not None and msg.get("id") == target:
            candidates.insert(0, text)
        else:
            candidates.append(text)
    for text in candidates:
        if parse_outline(text):
            return text
    return candidates[0] if candidates else ""


class Action:
    async def action(
        self,
        body: dict,
        __user__: dict = None,
        __event_emitter__=None,
        __request__=None,
    ):
        async def emit(event):
            if __event_emitter__:
                await __event_emitter__(event)

        async def status(desc, done=False):
            await emit({"type": "status", "data": {"description": desc, "done": done}})

        async def notify(content, kind="info"):
            # A model preset can hide the status line via its status_updates
            # capability, so every outcome also toasts and hits the server log.
            print(f"[openPPT] {kind}: {content}", flush=True)
            await emit({"type": "notification", "data": {"type": kind, "content": content}})

        try:
            messages = body.get("messages", [])
            target = body.get("id")
            content = _pick_outline(messages, target)
            slides = parse_outline(content)
            saved_chat = None
            if not slides:
                # The body's copy of the conversation is unreliable: some Open
                # WebUI versions POST the action assistant messages whose
                # 'content' is an empty string, so there is nothing to parse and
                # every export fails while the chat plainly holds an outline.
                # Fall back to the stored chat, which is also where the template
                # has had to come from since v0.5.1. Only on failure, so a body
                # that does carry the outline costs no extra lookup.
                saved_chat = await self._saved_chat(body)
                stored = _chat_messages(saved_chat)
                if stored:
                    from_chat = _pick_outline(stored, target)
                    slides = parse_outline(from_chat)
                    if slides or not content:
                        messages, content = stored, from_chat
            if not slides:
                # Temporary Chats are never written to the Chats table, so the
                # fallback above can never find one — and the live request body
                # is the only copy left, which different Open WebUI versions
                # populate with the assistant's text inconsistently (missing
                # message, empty string) since it's unsaved.
                # No stored copy at all — whether the lookup missed or there was
                # no chat_id to look up with. A missing chat_id is the strongest
                # unsaved-chat signal there is, so it must not suppress the hint.
                temp_chat = saved_chat is None
                # "Fix your outline" is the wrong advice when there is no reply
                # to have an outline in — and that's a distinct failure: neither
                # the request nor the saved chat holds an assistant message, so
                # the model's text never reached openPPT at all. Blaming the
                # outline sends people off rewriting a reply that was fine.
                no_reply = not any(
                    isinstance(m, dict) and m.get("role") == "assistant"
                    for m in list(messages or []) + _chat_messages(saved_chat)
                )
                if no_reply:
                    help_text = (
                        "openPPT never received the model's reply. Neither this "
                        "export request nor Open WebUI's saved copy of the chat "
                        "contains an assistant message, so there was no text to "
                        "read — this isn't about the outline's format. Reload the "
                        "page: if the reply disappears, Open WebUI isn't saving "
                        "it, so check that this isn't a Temporary Chat and that "
                        "chat history saving is on."
                    )
                else:
                    # Deliberately describes the grammar without writing a real
                    # tag: this text is posted back into the chat, and a literal
                    # slide element here would make the diagnostic itself parse
                    # as a deck on the next click (see the v0.3.5 guard in
                    # test_no_outline_diagnostic_is_not_itself_slide_shaped).
                    help_text = (
                        "openPPT: nothing slide-shaped in that message. Ask the model "
                        "for an HTML outline — one 'slide' element per slide, each "
                        "carrying a title attribute, with 'ul'/'li' bullets inside it."
                    )
                if temp_chat:
                    help_text += (
                        " This looks like a Temporary Chat — Open WebUI doesn't "
                        "save those, so openPPT has no reliable copy of the reply "
                        "to read. Turn off Temporary Chat and try again."
                    )
                # The toast/status is the only channel some deployments show —
                # the appended chat block below can go unseen — so lead with the
                # evidence, not the advice: the version proves the paste is
                # current, and 0 chars read (vs. N) separates "the action was
                # handed no text" from "the parser rejected real text". The
                # Temporary Chat note is the one exception worth repeating here.
                head = re.sub(r"\s+", " ", content[:80]).strip()
                temp_chat_note = (
                    "; looks like a Temporary Chat — turn it off and retry"
                    if temp_chat else ""
                )
                lead = (
                    "no assistant reply reached openPPT (not the outline's fault)"
                    if no_reply else "nothing slide-shaped here"
                )
                toast = (
                    f"openPPT {VERSION}: {lead} — read "
                    f"{len(content)} chars from {len(messages)} messages "
                    f"[{_shape(messages)}]; chat_id "
                    f"{'present' if body.get('chat_id') else 'MISSING'}; "
                    f"saved {_saved_shape(saved_chat)}; "
                    f"starts: {head or '(empty)'}{temp_chat_note}"
                )
                await notify(toast, "warning")
                await status(toast, done=True)
                # Toasts and the status line get truncated, so the detail goes in
                # the chat. Written so parse_outline finds nothing in it: no '---'
                # rule, no bullet list, no lone '**bold**' line — otherwise the
                # next click exports this diagnostic as a deck (v0.3.5 did).
                seen = content[:400] if content else "(nothing — the message was empty)"
                await emit(
                    {
                        "type": "message",
                        "data": {
                            "content": (
                                f"\n\n{APPENDIX_MARKER}\n"
                                f"**openPPT {VERSION} diagnostic** — {help_text} "
                                f"(messages in request: {len(messages)}, clicked id: "
                                f"`{target}`, read {len(content)} characters)\n\n"
                                f"What openPPT read from that message:\n\n"
                                f"```\n{seen}\n```\n"
                            )
                        },
                    }
                )
                return

            template = await self._find_template(body, __user__)
            wants_layout = any(s.get("layout") for s in slides)
            if wants_layout and template is None:
                await status(
                    "This deck names layouts but no .pptx/.potx template is "
                    "attached — using the default theme.",
                )

            note = " on your template" if template is not None else ""
            await status(f"Building PowerPoint ({len(slides)} slides{note})…")
            data = build_pptx(slides, template=template)

            # ponytail: internal Files/Storage API — swap to REST POST /api/v1/files/
            # with the user's token if these internals change across Open WebUI versions.
            from open_webui.models.files import FileForm, Files
            from open_webui.storage.provider import Storage

            slug = "".join(
                c if c.isalnum() else "-" for c in slides[0]["title"].lower()
            ).strip("-") or "deck"
            name = f"{slug}.pptx"
            file_id = str(uuid.uuid4())
            try:
                _, file_path = Storage.upload_file(
                    io.BytesIO(data), f"{file_id}_{name}", tags={}
                )
            except TypeError:  # older Open WebUI: no tags param
                _, file_path = Storage.upload_file(io.BytesIO(data), f"{file_id}_{name}")
            record = await _maybe_await(
                Files.insert_new_file(
                    __user__["id"],
                    FileForm(
                        id=file_id,
                        filename=name,
                        path=file_path,
                        data={},
                        meta={"name": name, "content_type": PPTX_MIME, "size": len(data)},
                    ),
                )
            )
            # insert_new_file returns None on failure instead of raising — without
            # this the link is posted and 404s ("download ready", nothing downloads).
            if record is None:
                raise RuntimeError("Open WebUI rejected the file record")

            # Two delivery paths on purpose. The 'files' event is how Open WebUI
            # returns its own generated files: it lands in the message's files
            # column and renders as an attachment chip (with a .pptx preview and
            # a download). The markdown link covers versions whose event handler
            # ignores 'files'. Emitted separately so one failing can't lose both.
            url = f"/api/v1/files/{file_id}/content"
            full_url = url
            if __request__ is not None:
                try:
                    full_url = str(__request__.base_url).rstrip("/") + url
                except Exception:
                    pass
            try:
                await emit(
                    {
                        "type": "files",
                        "data": {
                            "files": [
                                {
                                    "type": "file",
                                    "id": file_id,
                                    "url": file_id,  # FileItem prefixes /api/v1/files/
                                    "name": name,
                                    "collection_name": "",
                                    "status": "uploaded",
                                    "size": len(data),
                                    "content_type": PPTX_MIME,
                                    "meta": {
                                        "name": name,
                                        "content_type": PPTX_MIME,
                                        "size": len(data),
                                    },
                                }
                            ]
                        },
                    }
                )
            except Exception as e:
                print(f"[openPPT] files event rejected: {type(e).__name__}: {e}", flush=True)

            # The marker line lets a future export ignore everything below it
            # (see parse_outline) — otherwise re-clicking Export on this same
            # message would parse this link into a bogus trailing bullet. The
            # full URL repeated on its own line is copy-pasteable even where
            # neither the chip nor the markdown link renders.
            await emit(
                {
                    "type": "message",
                    "data": {
                        "content": (
                            f"\n\n{APPENDIX_MARKER}\n📊 [Download {name}]({full_url})"
                            f"\n\n{full_url}"
                        )
                    },
                }
            )
            await notify(f"{name} ready — {full_url}", "success")
            await status(f"{name} ready — {full_url}", done=True)
        except Exception as e:
            msg = f"openPPT export failed: {type(e).__name__}: {e}"
            await notify(msg, "error")
            await status(msg, done=True)

    async def _saved_chat(self, body: dict):
        """The stored chat dict behind this action, or None.

        Open WebUI strips the action request body down hard before POSTing to
        /api/chat/actions — it carries no 'files' at all, and on some versions
        the assistant message's 'content' comes through empty too. The stored
        chat is the only copy that reliably has both.
        """
        if not body.get("chat_id"):
            return None
        try:
            from open_webui.models.chats import Chats

            chat = await _maybe_await(Chats.get_chat_by_id(body["chat_id"]))
            return getattr(chat, "chat", None)
        except Exception as e:
            print(f"[openPPT] chat lookup failed: {type(e).__name__}: {e}", flush=True)
            return None

    async def _find_template(self, body: dict, __user__: dict):
        """Return template bytes from the most recent .pptx/.potx attached to
        the chat, or None.

        The action request body has no attachments in it: Open WebUI maps every
        message down to id/role/content/info/timestamp before POSTing to
        /api/chat/actions, and sends no top-level 'files'. So the attachment is
        read from the *saved* chat (body['chat_id']) — which keeps both the
        per-message 'files' and the chat-level 'files' list. The body is still
        scanned first, for versions or callers that do pass files along.
        """
        try:
            from open_webui.models.files import Files
            from open_webui.storage.provider import Storage
        except Exception:
            return None

        def file_ids(container):
            """(id, name) for each file in a {'messages': [...], 'files': [...]}
            dict — message attachments newest-first, then chat-level ones."""
            if not isinstance(container, dict):
                return
            for msg in reversed(_chat_messages(container)):
                if isinstance(msg, dict):
                    yield from _file_entries(msg.get("files"))
            yield from _file_entries(container.get("files"))

        sources = [body, await self._saved_chat(body)]

        for fid, hint in (e for src in sources for e in file_ids(src)):
            try:
                rec = await _maybe_await(Files.get_file_by_id(fid))
                if rec is None:
                    continue
                fname = (rec.filename or hint or "").lower()
                if not fname.endswith((".pptx", ".potx")):
                    continue
                with open(Storage.get_file(rec.path), "rb") as fh:
                    return fh.read()
            except Exception:
                continue
        return None
