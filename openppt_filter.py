"""
title: PowerPoint Template Primer
author: openPPT
version: 0.4.0
requirements: python-pptx
description: When you attach a .pptx/.potx template to the chat, this primes the model — before it answers — with your template's layout names and the openPPT HTML outline format, and tells it to distill the content you provided into a logically flowing deck that fills those layouts. Without a template, it still primes deck/presentation requests with the same outline format and a content-quality bar. Pair it with the "Export to PowerPoint" action: dump your content, and the model formats a deck you can export in one click. Works with any model (no tool calling needed).
"""

import io
import re

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pydantic import BaseModel, Field

# Marker so we never inject the primer twice into the same request.
SENTINEL = "[openPPT template primer]"

# Loose signal that the user is asking for a deck, used only when no
# template is attached (template presence is signal enough on its own).
_DECK_KEYWORDS_RE = re.compile(
    r"\b(deck|presentation|slides?|slideshow|powerpoint|pptx)\b", re.IGNORECASE
)

_QUALITY_BAR = """- Make every bullet concrete: one specific number, name, date, or action verb
  — never a topic label. If a bullet could describe any project unchanged
  ("Improved efficiency", "Key considerations"), replace it with the real
  fact from the content, or cut it.
- Give each slide exactly one point a viewer takes away, not a paragraph
  chopped into fragments. If two bullets are making the same point, merge
  them or cut one.
- Group by theme, not by the source's paragraph order — pull related facts
  together even if they were scattered across the input, and drop anything
  restating a point already made on an earlier slide.
- Size the deck to how much the content actually supports: a short input
  makes a short deck. Padding it out with filler slides ("Overview",
  "Conclusion") that add no new information is worse than a shorter deck.
- Push supporting detail, caveats, and sourcing into that slide's notes
  instead of cramming them onto the slide — the slide is the headline, the
  notes are the explanation."""

TITLE_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
SUBTITLE_TYPES = (PP_PLACEHOLDER.SUBTITLE,)
BODY_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)


def _ph_counts(layout):
    title = subtitle = body = 0
    for ph in layout.placeholders:
        try:
            t = ph.placeholder_format.type
        except Exception:
            t = None
        if t in TITLE_TYPES:
            title += 1
        elif t in SUBTITLE_TYPES:
            subtitle += 1
        elif t in BODY_TYPES:
            body += 1
    return title, subtitle, body


def _describe_layout(layout) -> str:
    """A short 'Name — what it's good for' hint from the layout's placeholders."""
    title, subtitle, body = _ph_counts(layout)
    if title and subtitle and not body:
        hint = "title slide (title + subtitle)"
    elif title and not subtitle and not body:
        hint = "divider / section break — title only, no bullets"
    elif title and body >= 3:
        hint = "comparison / multi-column"
    elif title and body == 2:
        hint = "two columns of bullets"
    elif title and body == 1:
        hint = "title + bullet list"
    elif not title and not body:
        hint = "blank"
    else:
        hint = "title + content"
    return f"- {layout.name} — {hint}"


def _layouts_from_bytes(data: bytes) -> list:
    prs = Presentation(io.BytesIO(data))
    return [_describe_layout(lay) for lay in prs.slide_layouts]


_FORMAT_RULES = """- Wrap every slide in <slide title="...">...</slide>. The first slide has
  only title and subtitle attributes, no body — that's the title slide.
- 3-6 <li> bullets per content slide, each under ~12 words. Nest a <ul> inside
  an <li> for a sub-point.
- Embed an image in place of a bullet: <img src="url" alt="description">.
- Put tabular data in <table><tr><td>...</td></tr></table>, not bullets.
- Put commands or code in <pre><code>...</code></pre>; it renders as a
  monospace box. Escape literal <, > and & in any text (as &lt; &gt; &amp;)
  so it can't be mistaken for markup.
- Respond with ONLY the <slide> elements — no prose before, after, or between
  them, and no ```html fence wrapping them."""


def _template_primer(layout_lines: list) -> str:
    layouts = "\n".join(layout_lines)
    return f"""{SENTINEL}
The user attached a PowerPoint template and provided content to turn into a
presentation. Distill the content they gave you into a deck — pull out the
substance, group it by theme, and cut anything that doesn't earn its place.
Do not invent facts that aren't in their content; you may condense and
rephrase.

Respond with ONLY an HTML outline in this exact format (no prose around it):

<slide title="Deck Title" subtitle="One-line subtitle"></slide>

<slide title="First Slide Title" layout="<layout name>">
<ul>
<li>Short bullet point</li>
<li>Another bullet point
<ul><li>Sub-point</li></ul>
</li>
</ul>
<notes>extra detail for the speaker (optional, kept off the slide)</notes>
</slide>

Formatting rules:
{_FORMAT_RULES}
- Open a new topic with a divider slide (layout with a title and nothing
  else), and lay out the deck so it flows: title → agenda/overview → grouped
  sections → summary/next steps.
{_QUALITY_BAR}

Choose a layout for each slide with the layout="<name>" attribute, using ONLY
these layouts from the attached template:
{layouts}

When the user later asks for changes, reply with the full revised outline."""


def _no_template_primer() -> str:
    return f"""{SENTINEL}
The user asked for a deck or presentation. Distill the content already in
this conversation into a deck — pull out the substance, group it by theme,
and cut anything that doesn't earn its place. Do not invent facts that
aren't there.

Respond with ONLY an HTML outline in this exact format (no prose around it):

<slide title="Deck Title" subtitle="One-line subtitle"></slide>

<slide title="First Slide Title">
<ul>
<li>Short bullet point</li>
<li>Another bullet point
<ul><li>Sub-point</li></ul>
</li>
</ul>
<notes>extra detail for the speaker (optional, kept off the slide)</notes>
</slide>

Formatting rules:
{_FORMAT_RULES}
- Flow: title -> agenda/overview -> grouped sections -> summary/next steps.
{_QUALITY_BAR}

If this isn't a deck/presentation request, ignore this and answer normally.
When the user later asks for changes, reply with the full revised outline."""


def _primer(layout_lines: list = None) -> str:
    return _template_primer(layout_lines) if layout_lines else _no_template_primer()


def _looks_like_deck_request(messages: list) -> bool:
    for m in reversed(messages or []):
        if m.get("role") == "user":
            return bool(_DECK_KEYWORDS_RE.search(m.get("content") or ""))
    return False


class Filter:
    class Valves(BaseModel):
        enabled: bool = Field(
            default=True,
            description="Prime the model whenever a .pptx/.potx template is attached.",
        )
        prime_without_template: bool = Field(
            default=True,
            description="Also prime deck/presentation requests when no template is attached, using a template-less primer.",
        )
        priority: int = Field(
            default=0, description="Filter execution priority (lower runs first)."
        )

    def __init__(self):
        self.valves = self.Valves()

    def inlet(self, body: dict, __user__: dict = None) -> dict:
        """Before the model answers: if a template is attached, inject a system
        message listing its layouts and the outline spec; otherwise, for a
        deck/presentation request, inject the template-less version of the
        same spec. Never raises — on any problem it returns the request
        untouched so the chat is unaffected."""
        try:
            if not self.valves.enabled:
                return body
            messages = body.get("messages", []) or []
            if any(
                m.get("role") == "system" and SENTINEL in (m.get("content") or "")
                for m in messages
            ):
                return body  # already primed this request

            data = self._find_template_bytes(body)
            if data is not None:
                layout_lines = _layouts_from_bytes(data)
                if not layout_lines:
                    return body
                content = _primer(layout_lines)
            elif self.valves.prime_without_template and _looks_like_deck_request(
                messages
            ):
                content = _primer()
            else:
                return body

            primer = {"role": "system", "content": content}
            # Sit right after an existing system prompt, else lead the messages.
            idx = 0
            for i, m in enumerate(messages):
                if m.get("role") == "system":
                    idx = i + 1
                else:
                    break
            messages.insert(idx, primer)
            body["messages"] = messages
        except Exception:
            return body
        return body

    def _find_template_bytes(self, body: dict):
        """Newest-first, return bytes of the first attached .pptx/.potx, or None.

        Duplicated from the export action on purpose: Open WebUI functions are
        pasted in individually and can't import one another.
        """
        try:
            from open_webui.models.files import Files
            from open_webui.storage.provider import Storage
        except Exception:
            return None

        def file_ids():
            for msg in reversed(body.get("messages", []) or []):
                for f in msg.get("files", []) or []:
                    fid = (f.get("file") or {}).get("id") or f.get("id")
                    if fid:
                        yield fid, (f.get("name") or (f.get("file") or {}).get("filename") or "")
            for f in body.get("files", []) or []:
                fid = (f.get("file") or {}).get("id") or f.get("id")
                if fid:
                    yield fid, (f.get("name") or (f.get("file") or {}).get("filename") or "")

        for fid, hint in file_ids():
            try:
                rec = Files.get_file_by_id(fid)
                if rec is None:
                    continue
                fname = (rec.filename or hint or "").lower()
                if not fname.endswith((".pptx", ".potx")):
                    continue
                with open(Storage.get_file(rec.path), "rb") as fh:
                    return fh.read()
            except Exception:
                continue
        return None
