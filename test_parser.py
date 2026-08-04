import io

from pptx import Presentation

from openppt_action import (
    APPENDIX_MARKER,
    VERSION,
    _find_layout_by_name,
    _pick_outline,
    build_pptx,
    list_layouts,
    parse_outline,
)

SAMPLE = """Here's your deck:

<slide title="Q3 Results" subtitle="Revenue up 12% YoY"></slide>

<slide title="Key Wins">
<ul>
<li>Launched EU region</li>
<li>Churn down to 2.1%
<ul><li>Best quarter ever</li></ul>
</li>
</ul>
<notes>mention the EU launch timeline if asked</notes>
</slide>

<slide title="Next Quarter">
<ol>
<li>Ship mobile app</li>
<li>Hire 2 SREs</li>
</ol>
</slide>

<slide title="Team Photo">
<ul><li><img src="https://example.com/team.jpg" alt="Team offsite"></li></ul>
</slide>
"""

TEMPLATED = """<slide title="Roadmap" layout="Section Header"></slide>

<slide title="Details" layout="Two Content">
<ul><li>Point one</li><li>Point two</li></ul>
</slide>
"""


def test_parse_outline():
    slides = parse_outline(SAMPLE)
    assert len(slides) == 4
    assert slides[0] == {
        "title": "Q3 Results",
        "subtitle": "Revenue up 12% YoY",
        "layout": "",
        "notes": "",
        "bullets": [],
        "table": [],
        "code": "",
    }
    assert slides[1]["bullets"] == [
        (0, "Launched EU region"),
        (0, "Churn down to 2.1%"),
        (1, "Best quarter ever"),
    ]
    assert slides[1]["notes"] == "mention the EU launch timeline if asked"
    assert slides[2]["bullets"] == [(0, "Ship mobile app"), (0, "Hire 2 SREs")]
    assert slides[3]["bullets"] == [
        (0, {"alt": "Team offsite", "image": "https://example.com/team.jpg"})
    ]
    assert parse_outline("no slide tags here") == []


def test_parse_layout_attribute():
    slides = parse_outline(TEMPLATED)
    assert slides[0]["title"] == "Roadmap"
    assert slides[0]["layout"] == "Section Header"
    assert slides[1]["title"] == "Details"
    assert slides[1]["layout"] == "Two Content"
    assert slides[1]["bullets"] == [(0, "Point one"), (0, "Point two")]


def test_subtitle_attribute_works_on_any_slide():
    slides = parse_outline(
        '<slide title="A" subtitle="first"></slide><slide title="B" subtitle="second"><li>x</li></slide>'
    )
    assert slides[0]["subtitle"] == "first"
    assert slides[1]["subtitle"] == "second"
    assert slides[1]["bullets"] == [(0, "x")]


def test_unclosed_slide_tags_are_tolerated():
    """A model that forgets '</slide>' shouldn't lose the next slide."""
    slides = parse_outline(
        '<slide title="One"><li>a<slide title="Two"><li>b</slide>'
    )
    assert [s["title"] for s in slides] == ["One", "Two"]
    assert slides[0]["bullets"] == [(0, "a")]
    assert slides[1]["bullets"] == [(0, "b")]


def test_bare_text_inside_slide_becomes_a_bullet():
    """Text not wrapped in '<p>'/'<li>' shouldn't be silently dropped."""
    (slide,) = parse_outline('<slide title="Bare">Just some text</slide>')
    assert slide["bullets"] == [(0, "Just some text")]


def test_p_tag_is_a_bullet():
    (slide,) = parse_outline(
        '<slide title="T"><p>First paragraph</p><p>Second paragraph</p></slide>'
    )
    assert slide["bullets"] == [(0, "First paragraph"), (0, "Second paragraph")]


def test_inline_tags_flatten_to_plain_text():
    (slide,) = parse_outline(
        '<slide title="T"><ul><li><b>bold</b> <code>code</code> '
        '<a href="http://x">link</a></li></ul></slide>'
    )
    assert slide["bullets"] == [(0, "bold code link")]


def test_notes_and_blockquote_both_feed_notes():
    (slide,) = parse_outline(
        '<slide title="T"><li>a</li><blockquote>quoted note</blockquote>'
        "<notes>and this</notes></slide>"
    )
    assert slide["notes"] == "quoted note\nand this"
    assert slide["bullets"] == [(0, "a")]


def test_deeply_nested_bullets_cap_at_level_4():
    nested = "<slide title=\"T\">"
    for _ in range(7):
        nested += "<ul><li>"
    nested += "deep"
    nested += "</li></ul>" * 7
    nested += "</slide>"
    (slide,) = parse_outline(nested)
    assert slide["bullets"] == [(4, "deep")]


def test_ignores_appended_download_link():
    """Re-exporting a message that already has a download link shouldn't
    turn that link into a bogus trailing bullet."""
    appended = (
        SAMPLE
        + f"\n{APPENDIX_MARKER}\n📊 [Download deck.pptx](/api/v1/files/abc/content)"
        + "\n\nhttp://host/api/v1/files/abc/content\n"
    )
    assert parse_outline(appended) == parse_outline(SAMPLE)


def test_appendix_marker_stops_even_with_an_unclosed_slide():
    """The marker has to win even when the model's own last slide was never
    closed — otherwise openPPT's own appended text gets folded into it."""
    unclosed = '<slide title="Real"><li>x' + f"\n\n{APPENDIX_MARKER}\n<slide title=\"stray\"><li>y</slide>"
    slides = parse_outline(unclosed)
    assert len(slides) == 1
    assert slides[0]["title"] == "Real"
    assert slides[0]["bullets"] == [(0, "x")]


def test_never_raises():
    junk = [
        "",
        "\x00�🙂" * 100,
        "<" * 500,
        "<li>" * 5000,
        "<pre>\nunclosed",
        "<slide title='" + "x" * 2000,
        "```\nunclosed fence",
    ]
    for j in junk:
        assert isinstance(parse_outline(j), list)


def test_pick_outline_falls_back_when_the_clicked_message_is_empty():
    """Open WebUI sometimes hands the action a message with no content at all;
    reporting 'no outline' at a chat that plainly has one is the failure that
    started the v0.3.5 diagnostic-exported-as-a-deck chain."""
    messages = [
        {"id": "a", "role": "assistant", "content": SAMPLE},
        {"id": "b", "role": "user", "content": "export it"},
        {"id": "c", "role": "assistant", "content": ""},
    ]
    assert _pick_outline(messages, "c") == SAMPLE
    assert _pick_outline(messages, None) == SAMPLE
    assert _pick_outline([], None) == ""
    assert _pick_outline([{"id": "c", "role": "assistant", "content": ""}], "c") == ""


def test_pick_outline_strips_what_openppt_appended():
    """Exporting twice must not turn our own download link into a bullet."""
    content = (
        SAMPLE
        + f"\n\n{APPENDIX_MARKER}\n\n"
        + "📊 [Download deck.pptx](/api/v1/files/x/content)\n"
    )
    slides = parse_outline(_pick_outline([{"id": "a", "role": "assistant", "content": content}], "a"))
    assert len(slides) == 4
    assert slides[-1]["bullets"] == [
        (0, {"alt": "Team offsite", "image": "https://example.com/team.jpg"})
    ]


def test_build_pptx():
    data = build_pptx(parse_outline(SAMPLE))
    assert data[:2] == b"PK"  # .pptx is a zip


def test_build_pptx_skips_unreachable_image():
    slides = parse_outline(
        '<slide title="Photo"><li><img src="https://nonexistent.invalid/x.png" alt="x"></li></slide>'
    )
    data = build_pptx(slides)
    assert data[:2] == b"PK"


def _make_template_bytes():
    """A 'template' is just an ordinary .pptx; build one and re-serialize it."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # a stray sample slide
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_build_with_template_clears_sample_slides():
    template = _make_template_bytes()
    slides = parse_outline(SAMPLE)
    data = build_pptx(slides, template=template)
    out = Presentation(io.BytesIO(data))
    # 4 outline slides, and the template's stray sample slide is gone.
    assert len(out.slides) == 4
    assert out.slides[0].shapes.title.text == "Q3 Results"


def test_layout_directive_selects_named_layout():
    template = _make_template_bytes()
    slides = parse_outline('<slide title="Roadmap" layout="Section Header"><li>item</li></slide>')
    data = build_pptx(slides, template=template)
    out = Presentation(io.BytesIO(data))
    assert out.slides[0].slide_layout.name == "Section Header"


def test_body_bullets_land_in_a_placeholder():
    slides = parse_outline('<slide title="Ideas"><ul><li>first</li><li>second</li></ul></slide>')
    data = build_pptx(slides)
    out = Presentation(io.BytesIO(data))
    texts = [
        p.text
        for shape in out.slides[0].placeholders
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
    ]
    assert "first" in texts and "second" in texts


def test_long_title_shrinks_instead_of_overlapping_the_body():
    """python-pptx can't recompute PowerPoint's autofit, so a long title keeps
    the template's ~44pt and spills across the slide over the bullets."""
    long_title = "openPPT: nothing slide-shaped in that message. " * 4
    out = Presentation(
        io.BytesIO(build_pptx(parse_outline(f'<slide title="{long_title}"><li>a</li></slide>')))
    )
    frame = out.slides[0].shapes.title.text_frame
    assert frame.word_wrap is True
    assert frame.paragraphs[0].font.size.pt <= 20


def test_short_title_keeps_the_templates_own_size():
    out = Presentation(
        io.BytesIO(build_pptx(parse_outline('<slide title="Q3 Results"><li>a</li></slide>')))
    )
    # None = inherited from the layout; we must not restyle a title that fits
    assert out.slides[0].shapes.title.text_frame.paragraphs[0].font.size is None


def test_dense_bullet_list_shrinks():
    items = "".join(f"<li>point {i}</li>" for i in range(14))
    out = Presentation(
        io.BytesIO(build_pptx(parse_outline(f'<slide title="Ideas"><ul>{items}</ul></slide>')))
    )
    (body,) = [p for p in out.slides[0].placeholders if p.placeholder_format.idx == 1]
    assert len(body.text_frame.paragraphs) == 14
    assert all(p.font.size.pt <= 14 for p in body.text_frame.paragraphs)


def test_find_layout_by_name_is_fuzzy():
    prs = Presentation()
    assert _find_layout_by_name(prs, "section header").name == "Section Header"
    assert _find_layout_by_name(prs, "TWO CONTENT").name == "Two Content"
    assert _find_layout_by_name(prs, "nope") is None


TABLE_OUTLINE = """<slide title="Q3 Numbers">
<p>revenue held</p>
<table>
<tr><th>Region</th><th>Rev</th><th>Growth</th></tr>
<tr><td>EMEA</td><td>4.2</td><td>12%</td></tr>
<tr><td>APAC</td><td>3.1</td><td>30%</td></tr>
</table>
</slide>
"""


def test_parse_table():
    (slide,) = parse_outline(TABLE_OUTLINE)
    assert slide["table"] == [
        ["Region", "Rev", "Growth"],
        ["EMEA", "4.2", "12%"],
        ["APAC", "3.1", "30%"],
    ]
    assert slide["bullets"] == [(0, "revenue held")]  # the rows aren't bullets


def test_build_renders_table():
    out = Presentation(io.BytesIO(build_pptx(parse_outline(TABLE_OUTLINE))))
    (table,) = [s.table for s in out.slides[0].shapes if s.has_table]
    assert len(table.rows) == 3 and len(table.columns) == 3
    assert table.cell(2, 0).text == "APAC"


def test_parse_pre_code_block():
    (slide,) = parse_outline(
        '<slide title="Deploy"><li>one step</li><pre><code>make ship\nssh box</code></pre></slide>'
    )
    assert slide["code"] == "make ship\nssh box"
    assert slide["bullets"] == [(0, "one step")]


def test_parse_fenced_code_block_inside_a_slide():
    """Models fall back to ``` out of habit even when told to write HTML."""
    (slide,) = parse_outline('<slide title="Deploy">\n```bash\nmake ship\nssh box\n```\n</slide>')
    assert slide["code"] == "make ship\nssh box"


def test_code_with_angle_brackets_does_not_confuse_the_parser():
    """A ``` fence is HTML-escaped before parsing, so code containing '<'/'>'
    (e.g. a real HTML/XML sample) can't be mistaken for slide markup."""
    (slide,) = parse_outline(
        '<slide title="XML">\n```xml\n<a href="x"><b>hi</b></a>\n```\n</slide>'
    )
    assert slide["code"] == '<a href="x"><b>hi</b></a>'
    assert slide["bullets"] == []


def test_build_renders_code_in_a_monospace_box():
    outline = '<slide title="Deploy"><pre><code>print(\'hi\')</code></pre></slide>'
    out = Presentation(io.BytesIO(build_pptx(parse_outline(outline))))
    boxes = [
        s for s in out.slides[0].shapes
        if s.has_text_frame and "print('hi')" in s.text_frame.text
    ]
    (box,) = boxes
    assert box.text_frame.paragraphs[0].font.name == "Consolas"
    # the empty bullet placeholder is gone, not left as a "click to add" box
    assert not [p for p in out.slides[0].placeholders if p.placeholder_format.idx == 1]


def test_outline_wrapped_in_a_fence_still_parses():
    """Models wrap their whole answer in ``` whatever the prompt says; that
    fence must not swallow the deck into one slide's code block."""
    slides = parse_outline(
        'Here you go:\n\n```html\n<slide title="One"><li>a</li></slide>\n'
        '<slide title="Two"><li>b</li></slide>\n```\n'
    )
    assert [s["title"] for s in slides] == ["One", "Two"]
    assert all(not s["code"] for s in slides)


def test_bullets_and_a_table_share_the_slide():
    out = Presentation(io.BytesIO(build_pptx(parse_outline(TABLE_OUTLINE))))
    slide = out.slides[0]
    (body,) = [p for p in slide.placeholders if p.placeholder_format.idx == 1]
    (table,) = [s for s in slide.shapes if s.has_table]
    assert "revenue held" in body.text_frame.text
    assert table.top >= body.top + body.height  # table sits below the bullets


def test_list_layouts_reports_names():
    names = [n for n, _ in list_layouts(None)]
    assert "Title Slide" in names and "Title and Content" in names


class _FakeRequest:
    base_url = "http://host.example/"


def _run_action(async_api: bool, content: str = SAMPLE, __request__=None,
                body_content=None, history=False):
    """Drive Action.action against stand-in Open WebUI Files/Storage modules.

    Open WebUI's Files API is sync in older versions and async in current ones,
    so both shapes have to work. The async fakes only record their effect
    *inside* the coroutine — a dropped 'await' therefore shows up as a missing
    file record rather than passing silently.
    """
    import asyncio
    import sys
    import tempfile
    import types

    from openppt_action import Action

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fh:
        fh.write(_make_template_bytes())
        template_path = fh.name

    inserted = {}

    class Record:
        filename = "brand.pptx"
        path = template_path

    class Files:
        @staticmethod
        def insert_new_file(user_id, form):
            async def _insert():
                inserted[form["id"]] = form
                return form

            if async_api:
                return _insert()
            inserted[form["id"]] = form
            return form

        @staticmethod
        def get_file_by_id(fid):
            async def _get():
                return Record()

            return _get() if async_api else Record()

    class Storage:
        @staticmethod
        def upload_file(stream, filename, tags=None):
            return stream.read(), f"/data/{filename}"

        @staticmethod
        def get_file(path):
            return path

    stored_messages = [
        {"id": "m0", "role": "user", "content": "deck please",
         "files": [{"type": "file", "id": "f1", "name": "brand.pptx"}]},
        {"id": "m1", "role": "assistant", "content": content},
    ]

    class Chat:
        # what Open WebUI persists: the attachment lives on the user message.
        # Some versions keep the flat list, others only the history id->msg map.
        chat = (
            {"history": {"messages": {m["id"]: m for m in stored_messages}}}
            if history
            else {"messages": stored_messages}
        )

    class Chats:
        @staticmethod
        def get_chat_by_id(cid):
            async def _get():
                return Chat()

            return _get() if async_api else Chat()

    files_mod = types.ModuleType("open_webui.models.files")
    files_mod.Files = Files
    files_mod.FileForm = lambda **kw: kw
    chats_mod = types.ModuleType("open_webui.models.chats")
    chats_mod.Chats = Chats
    storage_mod = types.ModuleType("open_webui.storage.provider")
    storage_mod.Storage = Storage
    fakes = {
        "open_webui": types.ModuleType("open_webui"),
        "open_webui.models": types.ModuleType("open_webui.models"),
        "open_webui.models.files": files_mod,
        "open_webui.models.chats": chats_mod,
        "open_webui.storage": types.ModuleType("open_webui.storage"),
        "open_webui.storage.provider": storage_mod,
    }

    events = []

    async def emitter(event):
        events.append(event)

    # Exactly what Open WebUI POSTs to /api/chat/actions: no 'files' anywhere,
    # each message stripped to id/role/content. The template has to be found
    # through chat_id, or every export silently falls back to the plain theme.
    body = {
        "id": "m1",
        "chat_id": "c1",
        "messages": [
            {"id": "m0", "role": "user", "content": "deck please"},
            {"id": "m1", "role": "assistant",
             "content": content if body_content is None else body_content},
        ],
    }
    sys.modules.update(fakes)
    try:
        asyncio.run(
            Action().action(
                body,
                __user__={"id": "u1"},
                __event_emitter__=emitter,
                __request__=__request__,
            )
        )
    finally:
        for key in fakes:
            sys.modules.pop(key, None)
    return events, inserted


def test_action_inserts_the_file_and_links_to_it():
    for async_api in (False, True):
        events, inserted = _run_action(async_api)
        descriptions = " ".join(
            str(e["data"].get("description", "")) + str(e["data"].get("content", ""))
            for e in events
        )
        assert "failed" not in descriptions, descriptions
        # the record has to exist, or the link 404s under a "ready" toast
        (file_id,) = inserted
        assert inserted[file_id]["meta"]["size"] > 0
        link = f"/api/v1/files/{file_id}/content"
        assert any(link in str(e["data"].get("content", "")) for e in events), events
        # and it's offered as a real attachment, not only as markdown
        (attached,) = [e for e in events if e["type"] == "files"]
        (chip,) = attached["data"]["files"]
        assert chip["id"] == file_id and chip["url"] == file_id
        assert chip["type"] == "file" and chip["name"].endswith(".pptx")
        # and the attached .pptx was picked up as the template
        assert "on your template" in descriptions


def test_no_outline_diagnostic_is_not_itself_slide_shaped():
    """v0.3.5 posted an error report written in '---' + '- ' Markdown, so the
    next click exported the error report as a deck — the diagnostic must
    still contain no '<slide>' tag under the new grammar."""
    events, inserted = _run_action(False, content="Sorry, I can't help with that.")
    assert inserted == {}  # nothing was built
    (message,) = [e for e in events if e["type"] == "message"]
    posted = message["data"]["content"]
    assert APPENDIX_MARKER in posted
    assert parse_outline(posted.split(APPENDIX_MARKER, 1)[1]) == []
    # and once appended to the message, it is invisible to the next click
    appended = {"id": "m1", "role": "assistant", "content": "Sorry, I can't help with that." + posted}
    assert _pick_outline([appended], "m1") == "Sorry, I can't help with that."


def test_no_outline_toast_carries_the_evidence():
    """Some deployments surface only the toast — the appended chat block goes
    unseen — so the toast itself has to say which paste ran and how much text
    the action was actually handed. '0 chars' (the action got nothing) and
    'N chars' (the parser rejected real text) are different root causes, and
    without them in the toast there is no way to tell them apart remotely."""
    for content, expect_chars in (("", 0), ("Sorry, I can't help with that.", 30)):
        events, inserted = _run_action(False, content=content)
        assert inserted == {}  # nothing was built
        toasts = [
            str(e["data"].get("content", "")) + str(e["data"].get("description", ""))
            for e in events
            if e["type"] in ("notification", "status")
        ]
        assert toasts, events
        blob = " ".join(toasts)
        assert VERSION in blob, blob  # a stale paste shows up here
        assert f"read {expect_chars} chars" in blob, blob
    # and the empty case names itself rather than trailing off into nothing
    assert "(empty)" in blob or "(empty)" in " ".join(
        str(e["data"].get("content", ""))
        for e in _run_action(False, content="")[0]
        if e["type"] == "notification"
    )


def test_outline_is_read_from_the_saved_chat_when_the_body_is_blank():
    """Some Open WebUI versions POST the action a body whose assistant message
    content is an empty string — the toast reported exactly
    'read 0 chars from 2 messages [user:str(12), assistant:str(0)]'.
    parse_outline is pure, so with no text there is nothing to parse and every
    export fails on that deployment while an identical chat exports fine on a
    version that populates the body. The outline has to come from the stored
    chat, the same way v0.5.1 had to go there for the template.
    """
    for hist in (False, True):
        events, inserted = _run_action(False, body_content="", history=hist)
        blob = " ".join(
            str(e["data"].get("description", "")) + str(e["data"].get("content", ""))
            for e in events
        )
        assert "nothing slide-shaped" not in blob, blob
        (file_id,) = inserted  # a real deck was built from the stored outline
        assert inserted[file_id]["meta"]["size"] > 0
        assert "on your template" in blob  # and the template still resolved

    # the body still wins when it does carry the outline: no needless chat read
    events, inserted = _run_action(False)
    assert inserted and "nothing slide-shaped" not in " ".join(
        str(e["data"].get("content", "")) for e in events
    )
    # and a genuinely outline-less chat still reports, rather than hanging on
    # a fallback that can't help
    events, inserted = _run_action(False, content="Sorry, I can't help.", body_content="")
    assert inserted == {}
    assert any("nothing slide-shaped" in str(e["data"].get("content", "")) for e in events)


def test_shape_separates_the_ways_content_goes_missing():
    """'read 0 chars' has four causes needing four different fixes. _shape has
    to tell them apart from the toast alone, with no access to the machine."""
    from openppt_action import _shape

    # list-shaped (multimodal) content — _pick_outline's isinstance guard
    # turns this into "" and the outline is lost
    assert "assistant:list(1)" in _shape(
        [{"role": "assistant", "content": [{"type": "text", "text": "<slide>"}]}]
    )
    # genuinely empty content in the POST body
    assert "assistant:str(0)" in _shape([{"role": "assistant", "content": ""}])
    # no assistant role at all — candidates never gets an entry
    assert _shape([{"role": "user", "content": "hi"}]) == "user:str(2)"
    # content absent entirely
    assert "assistant:NoneType(?)" in _shape([{"role": "assistant"}])
    # and it never raises on junk, since it runs inside the failure path
    assert _shape(None) == "(none)"
    assert _shape([]) == "(none)"
    assert "str!" in _shape(["not-a-dict"])


def test_action_appends_full_url_when_request_available():
    events, inserted = _run_action(False, __request__=_FakeRequest())
    (file_id,) = inserted
    full_link = f"http://host.example/api/v1/files/{file_id}/content"
    contents = [str(e["data"].get("content", "")) for e in events]
    (download_msg,) = [c for c in contents if "📊" in c]
    assert full_link in download_msg
    assert APPENDIX_MARKER in download_msg
    # the full URL is repeated as its own line, after everything else
    assert download_msg.rstrip().endswith(full_link)


if __name__ == "__main__":
    for _name, _fn in sorted(list(globals().items())):
        if _name.startswith("test_"):
            _fn()
    print("ok")
